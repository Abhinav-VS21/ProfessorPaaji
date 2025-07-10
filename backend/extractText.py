from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
import json
import sys


def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    slide_data = []

    for i, slide in enumerate(prs.slides):
        slide_dict = {
            "slide_id": i,
            "title": "",
            "body": "",
            "notes": "",
            "has_image": False,
        }
        max_title_length = 70
        min_title_length = 4

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide_dict["has_image"] = True

            if hasattr(shape, "text") and shape.text.strip():
                if (
                    shape.is_placeholder
                    and shape.placeholder_format == PP_PLACEHOLDER.TITLE
                ):
                    slide_dict["title"] = (
                        slide_dict["title"] + " " + str(shape.text.strip())
                    )
                else:
                    slide_dict["body"] = slide_dict["body"] + shape.text.strip()
        if min_title_length <= len(slide_dict["title"]) <= max_title_length:
            slide_dict["title"] = ""

        if not slide_dict["title"]:
            body_data = slide_dict["body"]
            endLimiter = ["\n", ".", ":", ";", "/"]
            positions = [
                body_data.find(e) for e in endLimiter if body_data.find(e) != -1
            ]
            endIndex = min(positions) if positions else -1

            if endIndex != -1 and min_title_length <= endIndex <= max_title_length:
                slide_dict["title"] = body_data[0:endIndex]
                slide_dict["body"] = body_data[endIndex:]

        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_dict["notes"] = notes_text

        slide_data.append(slide_dict)

    return slide_data


# Save to JSON file
def save_to_json(data, out_path):
    with open(out_path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2, ensure_ascii=False)


if __name__ == "__main__":

    if len(sys.argv) != 3:
        print("Usage: python extract_ppt_text.py input.pptx output.json")
        sys.exit(1)

    pptx_file = sys.argv[1]
    json_file = sys.argv[2]
    slides = extract_text_from_pptx(pptx_file)
    save_to_json(slides, json_file)
